from flask import Flask, request, after_this_request
from pptx import Presentation
from io import BytesIO
from flask_cors import CORS 
import requests
import json
import threading
import hashlib
from supabase import create_client, Client
import os
import time
from datetime import datetime

supabase_url = 'https://nwsyclpwsgqlpdwfmzku.supabase.co'
supabase_key = 'eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6Im53c3ljbHB3c2dxbHBkd2Ztemt1Iiwicm9sZSI6ImFub24iLCJpYXQiOjE3MjQ3NzYyMzEsImV4cCI6MjA0MDM1MjIzMX0.G88AjKBuUxGEsDeYMd10Ao1OylLW8op9WJlBTi4kJQ0'
supabase: Client = create_client(supabase_url, supabase_key)
now = datetime.now()

def delayed_cleanup(path1, path2, delay):
    time.sleep(delay)
    try:
        if os.path.exists(path1):
            os.remove(path1)
        if os.path.exists(path2):
            os.remove(path2)
    except Exception as e:
        print(f"Error deleting files: {e}")

def upload_file(file_path, bucket_name):
    file_name = os.path.basename(file_path)
    print(file_name)
    with open(file_path, 'rb') as file:
        response = supabase.storage.from_(bucket_name).upload(f'{file_name}', file)    
        if response.status_code != 200:
            print('Error uploading file:', response['error']['message'])
            return
        public_url = supabase.storage.from_(bucket_name).get_public_url(f'{file_name}')
        return public_url

def hash_md5(string):
    # Generate a random salt
    salt = os.urandom(16)


    return salt.hex()
app = Flask(__name__)
CORS(app) 

@app.route('/', methods=['POST'])
def createPPTFromOwnTemplates():
    url = request.json.get("url")
    jsons = request.json.get("json")
    file_name = hash_md5(f'{url}{now.strftime("%Y-%m-%d %H:%M:%S")}')
    destination = f"{file_name}.pptx"
    response = requests.get(url, stream=True)  
    if response.status_code == 200:
        with open(destination, 'wb') as file:
            for chunk in response.iter_content(chunk_size=8192):
                file.write(chunk)
        print(f"File downloaded successfully and saved to {destination}")
    else:
        print(f"Failed to download file. Status code: {response.status_code}")

    with open(destination, 'rb') as file:
        file_content = file.read()
    file = BytesIO(file_content)
    prs = Presentation(BytesIO(file.read()))
    
    # slides = {}
    # for slide_index, slide in enumerate(prs.slides):
    #     slides[f"slide{slide_index+1}"] = {}
    #     count_ind = 1
    #     for shape_index,shape in enumerate(slide.shapes):
    #         if shape.has_text_frame:
    #             text_frame = shape.text_frame
    #             for paragraph in text_frame.paragraphs:
    #                 for run in paragraph.runs:
    #                     txt = run.text.strip()
    #                     slides[f"slide{slide_index+1}"][f"current_text{count_ind}"] = txt
    #                     slides[f"slide{slide_index+1}"][f"updated_text{count_ind}"] = ""
    #                     count_ind = count_ind + 1
    # slides = json.dumps(slides, indent=4)
    # print(slides)
    # print(jsons)
    for slide_index, slide in enumerate(prs.slides):
        new_slide = jsons[f"slide{slide_index+1}"]
        count_ind = 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:                   
                    for run in paragraph.runs:
                        if new_slide[f'updated_text{count_ind}'] != "": 
                            print(new_slide[f'updated_text{count_ind}'])                        
                            run.text = run.text.replace(run.text.strip(),new_slide[f'updated_text{count_ind}'])
                        count_ind = count_ind + 1
    

    modified_ppt_path = f"modified_{destination}"
    prs.save(modified_ppt_path)
    print(f"Modified presentation saved to {modified_ppt_path}") 

    @after_this_request
    def cleanup(response):
       threading.Thread(target=delayed_cleanup, args=(modified_ppt_path,destination, 5)).start()
       return response    

    return upload_file(modified_ppt_path,"ppt")
if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
